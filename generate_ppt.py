import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Style Helpers: Premium Slate/Purple Theme
    DARK_BG = RGBColor(15, 23, 42)      # Slate-900 background
    PURPLE_TITLE = RGBColor(168, 85, 247)  # Purple title text
    TEXT_WHITE = RGBColor(255, 255, 255)
    TEXT_MUTED = RGBColor(148, 163, 184)  # Slate-400 description text

    blank_layout = prs.slide_layouts[6] # Blank slide layout

    # Slide 1: Title Slide
    slide1 = prs.slides.add_slide(blank_layout)
    slide1.background.fill.solid()
    slide1.background.fill.fore_color.rgb = DARK_BG

    title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.33), Inches(3.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    p1 = tf.paragraphs[0]
    p1.text = "SURAKSHA AI"
    p1.alignment = PP_ALIGN.CENTER
    p1.font.name = 'Trebuchet MS'
    p1.font.size = Pt(64)
    p1.font.bold = True
    p1.font.color.rgb = PURPLE_TITLE
    p1.space_after = Pt(10)

    p2 = tf.add_paragraph()
    p2.text = "National Disaster Warning Portal with Automated Regional Language Alert Routing"
    p2.alignment = PP_ALIGN.CENTER
    p2.font.name = 'Arial'
    p2.font.size = Pt(22)
    p2.font.color.rgb = TEXT_WHITE
    p2.space_after = Pt(20)

    p3 = tf.add_paragraph()
    p3.text = "Built with ReactJS, Tailwind CSS v4, Express, Node.js, and MongoDB Atlas"
    p3.alignment = PP_ALIGN.CENTER
    p3.font.name = 'Arial'
    p3.font.size = Pt(14)
    p3.font.italic = True
    p3.font.color.rgb = TEXT_MUTED

    # Slides 2-10 content
    slides_data = [
        {
            "title": "1. Project Vision & Problem Statement",
            "bullets": [
                "Language barriers pose a significant challenge in disaster alerts, where timely communication is a matter of life and death.",
                "Standard broadcasts in English or Hindi fail to warn rural populations in Southern India (Telangana, Tamil Nadu, Karnataka) effectively.",
                "Suraksha AI bridges this gap with an intelligent localized emergency warning and routing network.",
                "Provides crowdsourced disaster reporting alongside a controlled administrative portal for official review and alert broadcasting.",
                "Integrates public translation APIs and Twilio to immediately route localized warning messages in regional languages."
            ]
        },
        {
            "title": "2. Core Features: Public Dashboard",
            "bullets": [
                "Interactive Emergency Alert Feed: Displays only verified and published disaster bulletins categorized by severity levels.",
                "Incident Reporting System: Citizens can submit warnings indicating title, description, location details, state, and district.",
                "Optional Image Uploader: Allows attachable incident screenshots which compile into Base64 format and store in the database.",
                "Dynamic Region Dropdowns: State-to-District filtering populated using custom datasets mapped for Indian regions.",
                "SMS Alert Subscription: Public users can register their mobile numbers to automatically receive warning bulletins for their specific region."
            ]
        },
        {
            "title": "3. Core Features: Administrative Console",
            "bullets": [
                "Role-Based Access Control (RBAC): Restricted administrative login credentials to prevent unauthorized account creation.",
                "Manually Audited Approval Queue: Official review of crowdsourced alerts with options to 'Verify & Publish' or 'Reject'.",
                "Direct Alert Dispatcher: Custom broadcast panel to type alerts in English, filter target states/districts, and send warnings.",
                "Interactive Analytics Interface: Integrated visual charts monitoring severity distribution and subscriber volume statistics.",
                "Live Terminal Simulator: Built-in console showing detailed routing information (sender, target, state, translation comparison, status)."
            ]
        },
        {
            "title": "4. Technical Stack Architecture",
            "bullets": [
                "Frontend Architecture: React JS (single-page application) styled using Tailwind CSS v4 and Lucide-React icon packs.",
                "Data Visualizations: Recharts library rendering bar graphs and pie charts for real-time statistical overviews in the dashboard.",
                "Backend Framework: Node.js and Express REST API implementing route-level JWT (JSON Web Token) bearer authentication.",
                "Database Integration: MongoDB Atlas (hosted replica-set) handled securely via Mongoose Object Document Mapper (ODM).",
                "Integrations & Gateways: Libre/MyMemory API for real-time translation and Twilio API Client for SMS transmission."
            ]
        },
        {
            "title": "5. Dynamic Regional Routing Mappings",
            "bullets": [
                "Targeted Broadcasts: The system scans the database for subscriber numbers located specifically in the disaster's region.",
                "Automatic State-to-Language Mapping:",
                "  - Telangana Subscribers  -> Automatically mapped to Telugu (te)",
                "  - Tamil Nadu Subscribers  -> Automatically mapped to Tamil (ta)",
                "  - Karnataka Subscribers   -> Automatically mapped to Kannada (kn)",
                "  - Delhi & North India     -> Automatically mapped to Hindi (hi)",
                "This ensures alert relevance, reducing warning fatigue and dispatching warnings specifically to affected populations."
            ]
        },
        {
            "title": "6. AI-Powered Translation Integration",
            "bullets": [
                "Translation Engine: Integrates MyMemory Translation API to convert input English warnings into Telugu, Tamil, Kannada, or Hindi.",
                "Lightweight API Connection: Axios requests query the translator using dynamic language-pair query parameters (e.g. en|ta).",
                "Asynchronous Processing: Concurrent execution of translations for different recipient targets to optimize delivery speeds.",
                "Graceful Fallback: System defaults to the original English text in case of translation engine failure, ensuring alert delivery."
            ]
        },
        {
            "title": "7. Alert Dispatcher & Live Simulator Terminal",
            "bullets": [
                "The server checks for Twilio credentials in the environment configuration (.env) on startup.",
                "Production Mode: Sends physical SMS messages to verified regional subscribers via the Twilio SMS Gateway.",
                "Mock Simulation Mode: Gracefully falls back to mock logs if Twilio keys are missing, avoiding application failure.",
                "Broadcast Simulator Terminal: Interactive terminal in the dashboard showing original warning, translated text, target language, and recipient number."
            ]
        },
        {
            "title": "8. Seeding & Deployment Strategy",
            "bullets": [
                "Database Seeding: Automatically seeds a default administrator (admin@disaster.gov.in) and regional test numbers on start.",
                "Backend Hosting: REST API deployed to Render Cloud Platform (https://ai-disaster-management-system-c4xm.onrender.com).",
                "Frontend Connection: Client Axios instance configured with base URL targeting the production Render backend.",
                "State Management: React Context (AuthContext) persisting token data securely via localStorage."
            ]
        },
        {
            "title": "9. Project Impact & Future Scope",
            "bullets": [
                "Saves Lives: Removes critical language delays in communicating emergency situations to rural communities.",
                "Empowers Communities: Crowdsourced reports enable rapid emergency warnings before official teams arrive.",
                "Future Scope: Implement AI-powered image analysis to confirm hazard severity from photos.",
                "Interactive Voice Alerts (IVR): Add automated voice calls in regional languages for non-smartphone users.",
                "Offline Capabilities: Support mesh network Bluetooth warnings in zones with cellular network collapse."
            ]
        }
    ]

    for item in slides_data:
        slide = prs.slides.add_slide(blank_layout)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = DARK_BG

        # Title Box
        title_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.6), Inches(11.33), Inches(1.0))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        tf_title.margin_left = Inches(0)
        tf_title.margin_top = Inches(0)
        
        p = tf_title.paragraphs[0]
        p.text = item["title"]
        p.font.name = 'Trebuchet MS'
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = PURPLE_TITLE

        # Content Box
        content_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.33), Inches(5.0))
        tf_content = content_box.text_frame
        tf_content.word_wrap = True
        tf_content.margin_left = Inches(0)
        tf_content.margin_top = Inches(0)

        for i, bullet in enumerate(item["bullets"]):
            if i == 0:
                p_bullet = tf_content.paragraphs[0]
            else:
                p_bullet = tf_content.add_paragraph()
            
            p_bullet.text = bullet
            p_bullet.font.name = 'Arial'
            p_bullet.font.size = Pt(17)
            p_bullet.font.color.rgb = TEXT_WHITE
            p_bullet.space_after = Pt(12)
            
            # Add indentation for sub-items
            if bullet.strip().startswith("-"):
                p_bullet.level = 1
                p_bullet.font.size = Pt(16)
                p_bullet.font.color.rgb = TEXT_MUTED

    prs.save("Suraksha_AI_Disaster_Management_Presentation.pptx")
    print("Presentation created successfully: Suraksha_AI_Disaster_Management_Presentation.pptx")

if __name__ == '__main__':
    create_presentation()
